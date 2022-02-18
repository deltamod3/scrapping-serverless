import math
import operator
import urllib
import os
import pathlib
import re
import time
import urllib.request
import docx2txt
import numpy as np
import pandas
import textract
from PyPDF2 import PdfFileReader
from bs4 import BeautifulSoup
from odf.opendocument import load
from pptx import Presentation
from pyexcel_ods import get_data
from striprtf.striprtf import rtf_to_text
import nltk
import requests
from inscriptis import get_text
from nltk import WordNetLemmatizer
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from collections import Counter
from multiprocessing import Process, Pipe
from models import tmp_results_result_id, tmp_results_search_id, \
    get_dynamodb_table, tbl_tmp_searches, tmp_results_matched_similarity, tmp_results_date, \
    tmp_searches_description, tmp_searches_purpose, tmp_searches_technology, tmp_results_url, tbl_tmp_results_step1, \
    tbl_tmp_results_step2, tmp_results_title, tmp_results_description, tmp_results_text

file_dir = '/tmp'

nltk.data.path.append(file_dir)
nltk.download("punkt", download_dir=file_dir)
nltk.download("stopwords", download_dir=file_dir)
nltk.download("wordnet", download_dir=file_dir)


def handler(event, context):
    print('********************start here*****************')
    print(event)
    print('Records length:', len(event.get('Records')))

    for record in event.get('Records'):
        if record.get('eventName') == 'INSERT':
            result_id = int(record['dynamodb']['NewImage'][tmp_results_result_id]['N'])
            search_id = int(record['dynamodb']['NewImage'][tmp_results_search_id]['N'])
            tmp_results_step2 = get_dynamodb_table(tbl_tmp_results_step2)
            response = tmp_results_step2.get_item(
                Key={
                    tmp_results_search_id: search_id,
                    tmp_results_result_id: result_id
                })
            results_row = response.get('Item', None)
            if results_row:
                print('Existing item in tmp_results_step2 table: Skipped', results_row)
                continue
            # create a pipe for communication
            parent_conn, child_conn = Pipe()
            # create the process, pass instance and connection
            process = Process(target=process_event_record, args=(record, child_conn))
            process.start()
            child_conn.close()  # <-- Close the child_conn end in the main process too.
            try:
                event_id = parent_conn.recv()
                print('Processed this event:', event_id)
            except EOFError as err:
                print('Got error here: ', err)
            # parent_conn.close()
            process.join()
        elif record.get('eventName') == 'MODIFY':
            row = record['dynamodb']['NewImage']
            print('Modified row from tmp_results_step2', row)
        elif record.get('eventName') == 'DELETE':
            keys = record['dynamodb']['Keys']
            print('Removed a row from tmp_results_step2', keys)


def process_event_record(record, conn):
    result_id = int(record['dynamodb']['NewImage'][tmp_results_result_id]['N'])
    search_id = int(record['dynamodb']['NewImage'][tmp_results_search_id]['N'])
    url = record['dynamodb']['NewImage'][tmp_results_url]['S']
    title = record['dynamodb']['NewImage'][tmp_results_title]['S']
    description = record['dynamodb']['NewImage'][tmp_results_description]['S']

    try:
        print('Moving result_row: result_id: ', result_id)
        tmp_searches = get_dynamodb_table(tbl_tmp_searches)
        tmp_results_step1 = get_dynamodb_table(tbl_tmp_results_step1)
        tmp_results_step2 = get_dynamodb_table(tbl_tmp_results_step2)
        response = tmp_searches.get_item(Key={tmp_results_search_id: search_id})
        print('search row item: ', response)
        search_row = response.get('Item', None)
        if not response or not search_row:
            tmp_results_step1.delete_item(Key={
                tmp_results_result_id: result_id,
                tmp_results_search_id: search_id,
            })
            print('Deleted new item from tmp_results_step1: non-existing item in tmp_searches table', {
                tmp_results_result_id: result_id,
                tmp_results_search_id: search_id,
            })
            conn.send(record.get('eventID'))
            conn.close()
            return
        print('url:', url)
        url_type, text = parse_url(url)  # url to text, TODO Unused url_type for now
        print('parsed text:', text)
        technologies = search_row.get(tmp_searches_technology, None)
        if technologies and type(technologies) == list:
            different_text = []
            for tech in technologies:
                keywords = find_similarity_keyword(tech, text)  # technology filter
                ref_key = keywords["list_ref_key"]
                if ref_key["total_count"] != 0:
                    texts = {tech: text}
                    different_text.append(texts)
            print('technology_filter data:', different_text)
            if not different_text or len(different_text) == 0:
                tmp_results_step1.delete_item(Key={
                    tmp_results_result_id: result_id,
                    tmp_results_search_id: search_id,
                })
                print('Deleted new item from tmp_results_step1 with technology array but without ref_text: ', {
                    tmp_results_result_id: result_id,
                    tmp_results_search_id: search_id,
                })
                conn.send(record.get('eventID'))
                conn.close()
                return
        reference_text = search_row[tmp_searches_description]
        purpose = search_row[tmp_searches_purpose]
        rate = rate_searched_text(text, reference_text, purpose)  # Rate for text
        ts = int(time.time() * 1000)
        print('Current time: ', ts)
        new_item = {
            tmp_results_result_id: result_id,
            tmp_results_search_id: search_id,
            tmp_results_matched_similarity: rate['new_rating'],
            tmp_results_url: url,
            tmp_results_text: text,
            tmp_results_title: title,
            tmp_results_description: description,
            'url_type': url_type,
            tmp_results_date: ts,
        }
        tmp_results_step2.put_item(Item=new_item)
        print('Put new item into tmp_results_step2: ', new_item)
        tmp_results_step1.delete_item(Key={
            tmp_results_result_id: result_id,
            tmp_results_search_id: search_id,
        })
        print('Deleted new item from tmp_results_step1 put into tmp_results_step2: ', {
            tmp_results_result_id: result_id,
            tmp_results_search_id: search_id,
        })
    except Exception as e:
        print('Searching error: ', e)
    conn.send(record.get('eventID'))
    conn.close()


# url to text
def parse_url(url):
    url_text = pathlib.Path(url).suffix
    list_of_doc = ['.txt', '.PDF', '.pdf', '.csv', '.odp', '.ods', '.odt', '.rtf', '.docx', '.doc', '.xlsx', '.xls',
                   '.xlr', '.pptx', '.xml']
    if url_text not in list_of_doc:
        response = requests.get(url, "lxml", headers={'User-Agent': 'Mozilla/5.0'})
        text = get_text(response.text)
        remove_n = text.replace('\n', ' ')
        remove_r = remove_n.replace('\r', ' ')
        remove_t = remove_r.replace('\t', ' ')
        clean_text = re.sub(' +', ' ', remove_t)
        remove_leading_space = clean_text.lstrip()
        return 1, remove_leading_space
    else:
        if url.endswith('.txt'):
            data = urllib.request.urlopen(url)
            text = ''
            for line in data:
                decoded_line = line.decode("utf-8")
                text += decoded_line
            data = textclean(text)
            return 2, data
        elif url.endswith('.PDF') or url.endswith('.pdf'):
            r = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.pdf')
            open(file_path, 'wb').write(r.content)
            pdf_file = PdfFileReader(file_path)
            no_of_pages = pdf_file.getNumPages()
            pages_content = ''
            for i in range(no_of_pages):
                page = pdf_file.getPage(i)
                page_content = page.extractText()
                pages_content += page_content
            text = textclean(pages_content)
            return 2, text
        elif url.endswith('.csv'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.csv')
            open(file_path, 'wb').write(myfile.content)
            source = open(file_path, "r")
            decoded_file = source.read()
            text = textclean(decoded_file)
            return 2, text
        elif url.endswith('.odt'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.odt')
            open(file_path, 'wb').write(myfile.content)
            text = textract.process(file_path)
            esr = text.decode(encoding="utf-8")
            data = textclean(esr)
            return 2, data
        elif url.endswith('.odp'):
            from odf import text, teletype
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.csv')
            open(file_path, 'wb').write(myfile.content)
            textdoc = load(file_path)
            allparas = textdoc.getElementsByType(text.P)
            linelenght = len(allparas)
            texts = ''
            for line in range(linelenght):
                test = teletype.extractText(allparas[line])
                texts += test
            data = textclean(texts)
            return 2, data
        elif url.endswith('.ods'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.ods')
            open(file_path, 'wb').write(myfile.content)
            data1 = get_data(file_path)
            alphanumeric = [character for character in str(data1) if
                            character.isalnum() or character.isspace()]
            alphanumeric = "".join(alphanumeric)
            data = textclean(
                alphanumeric)
            return 2, data
        elif url.endswith('.docx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.docx')
            open(file_path, 'wb').write(myfile.content)
            text = docx2txt.process(file_path)
            data = textclean(text)
            return 2, data
        elif url.endswith('.doc'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            data = urllib.request.urlopen(url)
            file_path = os.path.join(file_dir, url_text + 'temp.doc')
            # with open(file_path, 'wb') as output:
            #     output.write(data.read())
            open(file_path, 'wb').write(myfile.content)
            text = textract.process(file_path)
            text = text.decode("utf-8")
            data = textclean(text)
            return 2, data
        elif url.endswith('.xlr'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.xlr')
            xlrfile = pandas.read_excel(file_path)
            json_str = xlrfile.to_json(orient='records')
            data = textclean(json_str)
            return 2, data
        elif url.endswith('.xlsx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.xlsx')
            open(file_path, 'wb').write(myfile.content)
            df = pandas.read_excel(file_path)
            df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
            data = df.to_json()
            return 2, data
        elif url.endswith('.xls'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.xls')
            open(file_path, 'wb').write(myfile.content)
            df = pandas.read_excel(file_path)
            df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
            os.remove(file_path)
            data = df.to_json()
            # data = self.textclean(df)
            return 2, data
        elif url.endswith('.xml'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            data = urllib.request.urlopen(url)
            file_path = os.path.join(file_dir, url_text + 'temp.xml')
            with open(file_path, 'wb') as output:
                output.write(data.read())
            with open(file_path, 'r') as f:
                data = f.read()
            bsdata = BeautifulSoup(data, "xml")
            test = str(bsdata)
            cleanr = re.compile('<.*?>')
            cleantext = re.sub(cleanr, '', test)
            data = textclean(cleantext)
            return 2, data
        elif url.endswith('.pptx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.pptx')
            open(file_path, 'wb').write(myfile.content)
            prs = Presentation(file_path)
            text_runs = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs += run.text
            data = textclean(text_runs)
            return 2, data
        elif url.endswith('.rtf'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, url_text + 'temp.rtf')
            open(file_path, 'wb').write(myfile.content)
            text = open(file_path, 'r')
            rtf = text.read()
            txt = rtf_to_text(rtf)
            data = textclean(txt)
            return 2, data


def textclean(text):
    bad_chars = ["*", '"', "'", '\\', '/']

    for i in bad_chars:
        text = text.replace(i, '')
    remove_n = text.replace('\n', ' ')
    remove_r = remove_n.replace('\r', ' ')
    remove_t = remove_r.replace('\t', ' ')
    clean_text = re.sub(' +', ' ', remove_t)
    return clean_text


# technology filter
def find_similarity_keyword(text, text2):
    wnl = WordNetLemmatizer()
    """For total Words of Test2"""
    split_list_text1 = text.split()
    split_list_text2 = text2.split()

    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text.lower())
    text2_tokens = word_tokenize(text2.lower())

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [wnl.lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [wnl.lemmatize(i) for i in tokens_without_sw_text2]

    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]
    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(
                    suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Creating "set" of both the keyword list"""
    setA = set(alphanumeric_for_text1)
    setB = set(alphanumeric_for_text2)

    """Separating the similar and non similar words From both the list of text"""
    overlap = setA & setB
    universe = setA | setB
    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)
    """Creating set For list of text1(FOR_: tot_unique_k)"""
    set_counter_text1 = set(counts_for_text1)

    """Counting repeated words(FOR_: tot_k_rep)"""
    count_rep_keyword_text1 = 0
    for key, count in counts_for_text1.items():
        count_rep_keyword_text1 += count

    """Calculating the repeated words in list of text2"""
    counts_for_text2 = Counter(alphanumeric_for_text2)
    set_counter = set(counts_for_text2)
    """Separating similar key_words from both the list of text"""
    similar_words = set_counter_text1 & set_counter
    tot_key_rep_text2 = 0
    tot_rep_key = 0
    """Count total repeated key_word from text1 to text2(FOR_: tot_ref_k) and total key_words(FOR_: repetition)"""
    for key, count in counts_for_text2.items():
        tot_key_rep_text2 += count
        if key in counts_for_text1:
            tot_rep_key += count  # + counts_for_text1[key]
    """Adding Total number of keywords and list of keywords"""
    unique_key_text1 = len(counts_for_text1)
    unique_key_text2 = len(counts_for_text2)
    list_keyword_text1 = counts_for_text1
    list_keyword_text2 = counts_for_text2

    """Adding Total number of keywords_repeated"""
    tot_key_rep_text1 = count_rep_keyword_text1
    tot_key_rep_text2 = tot_key_rep_text2

    """Adding Total number of keywords are present in text2 and list of those keywords"""
    tot_ref_key = len(overlap)
    list_ref_key = {}
    total_list_ref_count = 0
    for value in overlap:
        total_count = counts_for_text1[value] + counts_for_text2[value]
        list_ref_key[value] = total_count
        total_list_ref_count += total_count
    list_ref_key['total_count'] = total_list_ref_count

    """Set dictionary To descending order"""
    list_ref_key = dict(sorted(list_ref_key.items(), key=operator.itemgetter(1), reverse=True))
    list_keyword_text1 = dict(sorted(list_keyword_text1.items(), key=operator.itemgetter(1), reverse=True))
    list_keyword_text2 = dict(sorted(list_keyword_text2.items(), key=operator.itemgetter(1), reverse=True))

    """Calculating Percentage For both the list of Text"""
    result = float(len(overlap)) / len(universe) * 100
    data = {'percentage': result,
            'unique_key_text1': unique_key_text1,
            'unique_key_text2': unique_key_text2,
            'list_keyword_text1': list_keyword_text1,
            'list_keyword_text2': list_keyword_text2,
            'tot_key_rep_text1': tot_key_rep_text1,
            'tot_key_rep_text2': tot_key_rep_text2,
            'tot_ref_key': tot_ref_key,
            'list_ref_key': list_ref_key,
            'tot_rep_key': tot_rep_key,
            'total_word_text1': len(split_list_text1),
            'total_word_text2': len(split_list_text2)
            }
    return data


# Rate for text
def rate_searched_text(extracted_text, reference_text, purpose):
    extracted_text = extracted_text.lower()
    reference_text = reference_text.lower()
    ratings = replica_find_rating_percentage(extracted_text, reference_text)
    word_result = extracted_text
    if purpose:
        purpose = purpose.lower()
        word_result = replica_similarity_keyword(purpose, extracted_text)
    words = list(word_result)
    lines = find_unique_phrases(words, extracted_text)
    percentage = 0
    if purpose:
        test = replica_cosine_similarity(purpose, lines)
        percentage = test['percentage']
    purp = ratings + percentage
    new_rating = round(purp)
    rat = {"rating": ratings, "new_rating": new_rating, "purpose": percentage}
    print(rat)
    return rat


def find_keyword_repetition(text, text2):
    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text)
    text2_tokens = word_tokenize(text2)

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [WordNetLemmatizer().lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [WordNetLemmatizer().lemmatize(i) for i in tokens_without_sw_text2]
    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]

    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)

    """Calculating the repeated words in list of text2"""
    counts_for_text2 = Counter(alphanumeric_for_text2)

    tot_key_rep_text2 = 0
    tot_rep_key = 0
    """Count total repeated key_word from text1 to text2(FOR_: tot_ref_k) and total key_words(FOR_: repetition)"""
    for key, count in counts_for_text2.items():
        tot_key_rep_text2 += count
        if key in counts_for_text1:
            tot_rep_key += count  # + counts_for_text1[key]

    """Adding Total number of keywords_repeated"""
    tot_key_rep_text2 = tot_key_rep_text2

    """Calculating Percentage For both the list of Text"""
    result = tot_rep_key / tot_key_rep_text2 * 100

    data = {'percentage': result}
    return data


def text_to_cosine(text):
    search = str(text)
    WORD = re.compile(r'\w+')
    words = WORD.findall(search)
    return Counter(words)


def get_cosine(vec1, vec2):
    intersection = set(vec1.keys()) & set(vec2.keys())
    numerator = sum([vec1[x] * vec2[x] for x in intersection])
    sum1 = sum([vec1[x] ** 2 for x in vec1.keys()])
    sum2 = sum([vec2[x] ** 2 for x in vec2.keys()])
    denominator = math.sqrt(sum1) * math.sqrt(sum2)

    if not denominator or denominator == 0:
        return 0.0
    else:
        return float(numerator) / denominator


def find_cosine_similarity(text1, text2):
    vector1 = text_to_cosine(text1)
    vector2 = text_to_cosine(text2)
    cosine = get_cosine(vector1, vector2)
    print(cosine)
    Cosine_value = cosine * 100
    print(Cosine_value)
    return Cosine_value


def find_rating(cosine_value, keyword_repetition_value):
    if len(cosine_value) != len(keyword_repetition_value):
        return {"Error": 'Value is missing'}
    rating_result = []
    for index_value in range(len(cosine_value)):
        calculating_formula = cosine_value[index_value] * 75 + keyword_repetition_value[index_value] * 25
        calculating_parentage = calculating_formula / 100
        if keyword_repetition_value[index_value] >= 100:
            calculating_parentage = calculating_parentage + calculating_parentage * 25 / 100
        if calculating_parentage >= 100:
            calculating_parentage = 100
        rating_result.append({'rating_' + str(index_value): calculating_parentage})
    return {"Total_Rating": rating_result}


def find_cosine_similarity_with_keyword_repetition(text, text2):
    """Find Rating value"""

    """Take Result From cosine and keyword repetition"""
    cosine_value = find_cosine_similarity(text, text2)
    print(cosine_value)
    keyword_repetition_value = find_keyword_repetition(text, text2)
    print(keyword_repetition_value)
    """Find Rating value"""
    total_percentage = find_rating([cosine_value], [keyword_repetition_value['percentage']])
    print(total_percentage)
    return {'total_percentage': total_percentage['Total_Rating'][0]['rating_0'],
            'cosine_value': int(cosine_value),
            'keyword_repetition_value': int(keyword_repetition_value['percentage'])}


def replica_similarity_keyword(text, text2):
    """For total Words of Test2"""
    # split_list_text1 = text.split()
    # split_list_text2 = text2.split()

    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text)
    text2_tokens = word_tokenize(text2)

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [WordNetLemmatizer().lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [WordNetLemmatizer().lemmatize(i) for i in tokens_without_sw_text2]

    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]
    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(
                    suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Creating "set" of both the keyword list"""
    setA = set(alphanumeric_for_text1)
    setB = set(alphanumeric_for_text2)

    """Separating the similar and non similar words From both the list of text"""
    overlap = setA & setB
    universe = setA | setB
    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)
    """Creating set For list of text1(FOR_: tot_unique_k)"""
    set_counter_text1 = set(counts_for_text1)
    print(set_counter_text1)
    tot_ref_key = len(overlap)
    print(tot_ref_key)
    return set_counter_text1


def find_unique_phrases(words, text2):
    sentences = text2.split('.')
    lines = []
    for sentence in sentences:
        for word in words:
            if word in sentence:
                line = sentence
                lines.append(line)

    return lines


def replica_cosine_similarity(text1, text2):
    val = []
    new_lines = []
    phrases = []
    for line in text2:
        line = line.strip()
        if line not in new_lines:
            new_lines.append(line)
    lines = list(set(text2))
    for line in new_lines:
        vector1 = text_to_cosine(text1)
        vector2 = text_to_cosine(line)
        cosine = get_cosine(vector1, vector2)
        Cosine_value = cosine * 100
        if Cosine_value != 0.0:
            value = {Cosine_value}
            linecoisne = {line: value}
            phrases.append(linecoisne)
            val.append(value)
    number = len(val)
    percentage = []
    for vals in val:
        word_val = sum(vals)
        percentage.append(word_val)
    results = 0
    if number != 0:
        results = sum(percentage) / number
    data = {"percentage": results, "phrases": phrases}
    return data


def replica_find_rating_percentage(extracted_text, reference_text):
    similarity_dict = find_cosine_similarity_with_keyword_repetition(reference_text, str(extracted_text))
    similarity = similarity_dict['total_percentage']
    return similarity
