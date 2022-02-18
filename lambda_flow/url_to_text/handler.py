import io
import os
import pathlib
import re
import subprocess
import time
import urllib.request
from multiprocessing import Process, Pipe
import docx2txt
import numpy as np
import pandas
import requests
import textract
from PyPDF2 import PdfFileReader
from bs4 import BeautifulSoup
from inscriptis import get_text
from odf.opendocument import load
from pptx import Presentation
from pyexcel_ods import get_data
from striprtf.striprtf import rtf_to_text
from models import tmp_results_result_id, tmp_results_search_id, tmp_results_url, get_dynamodb_table, tbl_tmp_results_step1, tbl_tmp_results_step2, tmp_results_date, \
    tmp_results_text, \
    tmp_results_type, tbl_tmp_searches

file_dir = '/tmp'


def handler(event, context):
    print(event)
    print('Records length:', len(event.get('Records')))

    for record in event.get('Records'):
        # create a pipe for communication
        parent_conn, child_conn = Pipe()
        # create the process, pass instance and connection
        process = Process(target=process_event_record, args=(record, child_conn))
        process.start()
        child_conn.close()  # <-- Close the child_conn end in the main process too.
        try:
            event_id = parent_conn.recv()
            print('Processed this event:', event_id)
        except EOFError as err:
            print('Got error here: ', err)
        process.join()


def process_event_record(record, conn):
    if record.get('eventName') == 'INSERT':
        result_id = int(record['dynamodb']['NewImage'][tmp_results_result_id]['N'])
        search_id = int(record['dynamodb']['NewImage'][tmp_results_search_id]['N'])
        url = record['dynamodb']['NewImage'][tmp_results_url]['S']

        try:
            tmp_searches = get_dynamodb_table(tbl_tmp_searches)
            tmp_results_step1 = get_dynamodb_table(tbl_tmp_results_step1)
            response = tmp_searches.get_item(Key={tmp_results_search_id: search_id})
            print('search row response: ', response)
            search_row = response.get('Item', None)
            if not response or not search_row:
                tmp_results_step1.delete_item(Key={
                    tmp_results_result_id: result_id,
                    tmp_results_search_id: search_id,
                })
                print('Deleted new item from tmp_results_step1: non-existing item in tmp_searches table', {
                    tmp_results_result_id: result_id,
                    tmp_results_search_id: search_id,
                })
                conn.send(record.get('eventID'))
                conn.close()
                return
            tmp_results_step2 = get_dynamodb_table(tbl_tmp_results_step2)
            url_type, text = parse_url(url)
            ts = int(time.time() * 1000)
            print('Current time: ', ts)
            new_item = {
                tmp_results_result_id: result_id,
                tmp_results_search_id: search_id,
                tmp_results_text: text,
                tmp_results_type: url_type,
                tmp_results_date: ts,
            }
            tmp_results_step2.put_item(Item=new_item)
            print('Put new item into tmp_results_step2: ', new_item)
            tmp_results_step1.delete_item(Key={
                tmp_results_result_id: result_id,
                tmp_results_search_id: search_id,
            })
            print('Deleted new item from tmp_results_step1: ', {
                tmp_results_result_id: result_id,
                tmp_results_search_id: search_id,
            })
        except Exception as e:
            print('Searching error: ', e)
    elif record.get('eventName') == 'MODIFY':
        row = record['dynamodb']['NewImage']
        print('Modified row from tmp_results_step1', row)
    elif record.get('eventName') == 'DELETE':
        keys = record['dynamodb']['Keys']
        print('Removed a row from tmp_results_step1:', keys)
    conn.send(record.get('eventID'))
    conn.close()


def parse_url(url):
    url_text = pathlib.Path(url).suffix
    list_of_doc = ['.txt', '.PDF', '.pdf', '.csv', '.odp', '.ods', '.odt', '.rtf', '.docx', '.doc', '.xlsx', '.xls',
                   '.xlr', '.pptx', '.xml']
    if url_text not in list_of_doc:
        response = requests.get(url, "lxml", headers={'User-Agent': 'Mozilla/5.0'})
        text = get_text(response.text)
        remove_n = text.replace('\n', ' ')
        remove_r = remove_n.replace('\r', ' ')
        remove_t = remove_r.replace('\t', ' ')
        clean_text = re.sub(' +', ' ', remove_t)
        remove_leading_space = clean_text.lstrip()
        return 1, remove_leading_space
    else:
        if url.endswith('.txt'):
            data = urllib.request.urlopen(url)
            text = ''
            for line in data:
                decoded_line = line.decode("utf-8")
                text += decoded_line
            data = textclean(text)
            return 2, data
        elif url.endswith('.PDF') or url.endswith('.pdf'):
            r = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.pdf')
            open(file_path, 'wb').write(r.content)
            pdf_file = PdfFileReader(file_path)
            no_of_pages = pdf_file.getNumPages()
            pages_content = ''
            for i in range(no_of_pages):
                page = pdf_file.getPage(i)
                page_content = page.extractText()
                pages_content += page_content
            text = textclean(pages_content)
            return 2, text
        elif url.endswith('.csv'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.csv')
            open(file_path, 'wb').write(myfile.content)
            source = open(file_path, "r")
            decoded_file = source.read()
            text = textclean(decoded_file)
            return 2, text
        elif url.endswith('.odt'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.odt')
            open(file_path, 'wb').write(myfile.content)
            text = textract.process(file_path)
            esr = text.decode(encoding="utf-8")
            data = textclean(esr)
            return 2, data
        elif url.endswith('.odp'):
            from odf import text, teletype
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.csv')
            open(file_path, 'wb').write(myfile.content)
            textdoc = load(file_path)
            allparas = textdoc.getElementsByType(text.P)
            linelenght = len(allparas)
            texts = ''
            for line in range(linelenght):
                test = teletype.extractText(allparas[line])
                texts += test
            data = textclean(texts)
            return 2, data
        elif url.endswith('.ods'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.ods')
            open(file_path, 'wb').write(myfile.content)
            data1 = get_data(file_path)
            alphanumeric = [character for character in str(data1) if
                            character.isalnum() or character.isspace()]
            alphanumeric = "".join(alphanumeric)
            data = textclean(
                alphanumeric)
            return 2, data
        elif url.endswith('.docx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.docx')
            open(file_path, 'wb').write(myfile.content)
            text = docx2txt.process(file_path)
            data = textclean(text)
            return 2, data
        elif url.endswith('.doc'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            data = urllib.request.urlopen(url)
            file_path = os.path.join(file_dir, 'temp.doc')
            # with open(file_path, 'wb') as output:
            #     output.write(data.read())
            open(file_path, 'wb').write(myfile.content)
            text = textract.process(file_path)
            text = text.decode("utf-8")
            data = textclean(text)
            return 2, data
        elif url.endswith('.xlr'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.xlr')
            xlrfile = pandas.read_excel(file_path)
            json_str = xlrfile.to_json(orient='records')
            data = textclean(json_str)
            return 2, data
        elif url.endswith('.xlsx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.xlsx')
            open(file_path, 'wb').write(myfile.content)
            df = pandas.read_excel(file_path)
            df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
            data = df.to_json()
            return 2, data
        elif url.endswith('.xls'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.xls')
            open(file_path, 'wb').write(myfile.content)
            df = pandas.read_excel(file_path)
            df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
            os.remove(file_path)
            data = df.to_json()
            # data = self.textclean(df)
            return 2, data
        elif url.endswith('.xml'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            data = urllib.request.urlopen(url)
            file_path = os.path.join(file_dir, 'temp.xml')
            with open(file_path, 'wb') as output:
                output.write(data.read())
            with open(file_path, 'r') as f:
                data = f.read()
            bsdata = BeautifulSoup(data, "xml")
            test = str(bsdata)
            cleanr = re.compile('<.*?>')
            cleantext = re.sub(cleanr, '', test)
            data = textclean(cleantext)
            return 2, data
        elif url.endswith('.pptx'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.pptx')
            open(file_path, 'wb').write(myfile.content)
            prs = Presentation(file_path)
            text_runs = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs += run.text
            data = textclean(text_runs)
            return 2, data
        elif url.endswith('.rtf'):
            myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
            file_path = os.path.join(file_dir, 'temp.rtf')
            open(file_path, 'wb').write(myfile.content)
            text = open(file_path, 'r')
            rtf = text.read()
            txt = rtf_to_text(rtf)
            data = textclean(txt)
            return 2, data


def textclean(text):
    bad_chars = ["*", '"', "'", '\\', '/']

    for i in bad_chars:
        text = text.replace(i, '')
    remove_n = text.replace('\n', ' ')
    remove_r = remove_n.replace('\r', ' ')
    remove_t = remove_r.replace('\t', ' ')
    clean_text = re.sub(' +', ' ', remove_t)
    return clean_text


# PDF parser in progress
def decompress_pdf(temp_buffer):
    print('>>>>>>>>>>>>>>>111111', temp_buffer)
    fileno = temp_buffer.fileno()
    print('>>>>>>>>>>>>>>>111111', fileno)
    print('>>>>>>>>>>>>>>>111111', temp_buffer.seek(0))
    print('>>>>>>>>>>>>>>>111111', temp_buffer.seekable())
    temp_buffer.seek(0)  # Make sure we're at the start of the file.

    process = subprocess.Popen(['pdftk.exe',
                                '-',  # Read from stdin.
                                'output',
                                '-',  # Write to stdout.
                                'uncompress'],
                               stdin=temp_buffer,
                               stdout=subprocess.PIPE,
                               stderr=subprocess.PIPE)
    stdout, stderr = process.communicate()

    return io.StringIO(stdout)
