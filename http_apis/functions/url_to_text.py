import io
import os
import pathlib
import re
import subprocess
import urllib.request
import docx2txt
import numpy as np
import pandas
import requests
import textract
from PyPDF2 import PdfFileReader
from bs4 import BeautifulSoup
from inscriptis import get_text
from odf.opendocument import load
from pptx import Presentation
from pyexcel_ods import get_data
from striprtf.striprtf import rtf_to_text
from response_lib.helper import api_response, internal_server_error, bad_request_error

file_dir = '/mnt/efs'

if not os.path.isdir(file_dir):
    os.makedirs(file_dir)


def handler(event, context):
    print(event)
    request_qs = event.get('queryStringParameters', None)
    if not request_qs:
        return bad_request_error("Missing query string parameter")
    url = request_qs.get('url', None)
    if not url:
        return bad_request_error("Missing url")
    url_text = pathlib.Path(url).suffix
    list_of_doc = ['.txt', '.PDF', '.pdf', '.csv', '.odp', '.ods', '.odt', '.rtf', '.docx', '.doc', '.xlsx', '.xls',
                   '.xlr', '.pptx', '.xml']
    if url_text not in list_of_doc:
        try:
            response = requests.get(url, "lxml", headers={'User-Agent': 'Mozilla/5.0'})
            text = get_text(response.text)
            remove_n = text.replace('\n', ' ')
            remove_r = remove_n.replace('\r', ' ')
            remove_t = remove_r.replace('\t', ' ')
            clean_text = re.sub(' +', ' ', remove_t)
            remove_leading_space = clean_text.lstrip()
            return api_response(remove_leading_space)
        except Exception as e:
            return internal_server_error(e)
    else:
        if url.endswith('.txt'):
            try:
                data = urllib.request.urlopen(url)
                text = ''
                for line in data:
                    decoded_line = line.decode("utf-8")
                    text += decoded_line
                data = textclean(text)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.PDF') or url.endswith('.pdf'):
            try:
                r = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.pdf')
                open(file_path, 'wb').write(r.content)
                pdf_file = PdfFileReader(file_path)
                no_of_pages = pdf_file.getNumPages()
                pages_content = ''
                for i in range(no_of_pages):
                    page = pdf_file.getPage(i)
                    page_content = page.extractText()
                    pages_content += page_content
                text = textclean(pages_content)
                return api_response(text)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.csv'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.csv')
                open(file_path, 'wb').write(myfile.content)
                source = open(file_path, "r")
                decoded_file = source.read()
                text = textclean(decoded_file)
                return api_response(text)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.odt'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.odt')
                open(file_path, 'wb').write(myfile.content)
                text = textract.process(file_path)
                esr = text.decode(encoding="utf-8")
                data = textclean(esr)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.odp'):
            try:
                from odf import text, teletype
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.csv')
                open(file_path, 'wb').write(myfile.content)
                textdoc = load(file_path)
                allparas = textdoc.getElementsByType(text.P)
                linelenght = len(allparas)
                texts = ''
                for line in range(linelenght):
                    test = teletype.extractText(allparas[line])
                    texts += test
                data = textclean(texts)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.ods'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.ods')
                open(file_path, 'wb').write(myfile.content)
                data1 = get_data(file_path)
                alphanumeric = [character for character in str(data1) if
                                character.isalnum() or character.isspace()]
                alphanumeric = "".join(alphanumeric)
                data = textclean(
                    alphanumeric)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.docx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.docx')
                open(file_path, 'wb').write(myfile.content)
                text = docx2txt.process(file_path)
                data = textclean(text)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.doc'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                data = urllib.request.urlopen(url)
                file_path = os.path.join(file_dir, 'temp.doc')
                # with open(file_path, 'wb') as output:
                #     output.write(data.read())
                open(file_path, 'wb').write(myfile.content)
                text = textract.process(file_path)
                text = text.decode("utf-8")
                data = textclean(text)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.xlr'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.xlr')
                xlrfile = pandas.read_excel(file_path)
                json_str = xlrfile.to_json(orient='records')
                data = textclean(json_str)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.xlsx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.xlsx')
                open(file_path, 'wb').write(myfile.content)
                df = pandas.read_excel(file_path)
                df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                df = df.to_json()
                return api_response(df)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.xls'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.xls')
                open(file_path, 'wb').write(myfile.content)
                df = pandas.read_excel(file_path)
                df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                os.remove(file_path)
                df = df.to_json()
                # data = self.textclean(df)
                return api_response(df)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.xml'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                data = urllib.request.urlopen(url)
                file_path = os.path.join(file_dir, 'temp.xml')
                with open(file_path, 'wb') as output:
                    output.write(data.read())
                with open(file_path, 'r') as f:
                    data = f.read()
                bsdata = BeautifulSoup(data, "xml")
                test = str(bsdata)
                cleanr = re.compile('<.*?>')
                cleantext = re.sub(cleanr, '', test)
                data = textclean(cleantext)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.pptx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.pptx')
                open(file_path, 'wb').write(myfile.content)
                prs = Presentation(file_path)
                text_runs = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs += run.text
                data = textclean(text_runs)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)
        elif url.endswith('.rtf'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(file_dir, 'temp.rtf')
                open(file_path, 'wb').write(myfile.content)
                text = open(file_path, 'r')
                rtf = text.read()
                txt = rtf_to_text(rtf)
                data = textclean(txt)
                return api_response(data)
            except Exception as e:
                return internal_server_error(e)


def textclean(text):
    bad_chars = ["*", '"', "'", '\\', '/']

    for i in bad_chars:
        text = text.replace(i, '')
    remove_n = text.replace('\n', ' ')
    remove_r = remove_n.replace('\r', ' ')
    remove_t = remove_r.replace('\t', ' ')
    clean_text = re.sub(' +', ' ', remove_t)
    return clean_text


# PDF parser in progress
def decompress_pdf(temp_buffer):
    print('>>>>>>>>>>>>>>>111111', temp_buffer)
    fileno = temp_buffer.fileno()
    print('>>>>>>>>>>>>>>>111111', fileno)
    print('>>>>>>>>>>>>>>>111111', temp_buffer.seek(0))
    print('>>>>>>>>>>>>>>>111111', temp_buffer.seekable())
    temp_buffer.seek(0)  # Make sure we're at the start of the file.

    process = subprocess.Popen(['pdftk.exe',
                                '-',  # Read from stdin.
                                'output',
                                '-',  # Write to stdout.
                                'uncompress'],
                               stdin=temp_buffer,
                               stdout=subprocess.PIPE,
                               stderr=subprocess.PIPE)
    stdout, stderr = process.communicate()

    return io.StringIO(stdout)
